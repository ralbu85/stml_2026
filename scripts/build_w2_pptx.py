"""Build lectures/week02/slides-draft-w2.pptx from slides-content-w2.md in the existing deck's format
(10 x 5.625 in, white background, Arial, 24pt title / 16pt grey bullets, 36pt dividers, 44/20pt covers).
Figures (PNG renders of figures/*.svg), a real table on the method-family slide, monospace prompt quotes,
speaker notes from <!-- --> comments."""
import re
import sys
import pathlib

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = pathlib.Path(sys.argv[1])
FIGDIR = pathlib.Path(sys.argv[2])
OUT = pathlib.Path(sys.argv[3])

BLACK = RGBColor(0x00, 0x00, 0x00)
BODY = RGBColor(0x44, 0x44, 0x44)
GREY = RGBColor(0x88, 0x88, 0x88)
QUOTE = RGBColor(0x1F, 0x3A, 0x5F)
FONT = "Arial"
MONO = "Courier New"

# ---------- parse the draft ----------
text = SRC.read_text(encoding="utf-8")
body = text.split("\n---\n", 1)[1]
slides, cur = [], None
for raw in body.split("\n"):
    line = raw.rstrip()
    m = re.match(r"^## (\d+)\. (.*)$", line)
    if m:
        cur = {"n": int(m.group(1)), "title": m.group(2).strip(), "items": [], "notes": []}
        slides.append(cur)
        continue
    if cur is None or line.strip() == "---":
        continue
    cm = re.match(r"^<!--\s*(.*?)\s*-->$", line)
    if cm:
        cur["notes"].append(cm.group(1))
    elif line.startswith("  - "):
        cur["items"].append(("sub", line[4:]))
    elif line.startswith("- "):
        cur["items"].append(("bullet", line[2:]))
    elif line.startswith("> "):
        cur["items"].append(("quote", line[2:]))
    elif line.strip() == ">":
        cur["items"].append(("quote", ""))
    elif line.startswith("|"):
        cur["items"].append(("table", line))
    elif line.startswith("[Figure"):
        cur["items"].append(("figure", re.search(r"`figures/(.*?)`", line).group(1)))
    elif line.startswith("Caption:"):
        cur["items"].append(("caption", line[len("Caption:"):].strip()))
    elif line.strip():
        cur["items"].append(("text", line.strip()))

# ---------- deck geometry (from the committed draft) ----------
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(9144000), Emu(5143500)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TITLE_BOX = (Emu(384048), Emu(256032), Emu(8375904), Emu(640080))
BODY_BOX = (Emu(502920), Emu(987552), Emu(8229600), Emu(3931920))


def strip_md(s):
    return re.sub(r"`(.*?)`", r"\1", s)


def run(par, txt, size, color=BODY, bold=False, font=FONT):
    r = par.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.name = font
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def textbox(slide, box):
    tb = slide.shapes.add_textbox(*box)
    tb.text_frame.word_wrap = True
    return tb


def cover(slide, title, sub):
    tb = textbox(slide, (Emu(548640), Emu(1737360), Emu(8046720), Emu(1645920)))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, title, 44, BLACK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run(p2, sub, 20, GREY)


def divider(slide, title):
    tb = textbox(slide, (Emu(548640), Emu(1828800), Emu(8046720), Emu(1463040)))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, title, 36, BLACK)


def title(slide, txt):
    tb = textbox(slide, TITLE_BOX)
    run(tb.text_frame.paragraphs[0], txt, 24, BLACK)


def body_size(items):
    chars = sum(len(t) for k, t in items if k in ("bullet", "sub", "text", "quote", "caption"))
    if chars > 1000:
        return 13
    if chars > 780:
        return 14
    if chars > 620:
        return 15
    return 16


def add_body(slide, items, box, size):
    tb = textbox(slide, box)
    tf = tb.text_frame
    first = True
    for kind, txt in items:
        if kind not in ("bullet", "sub", "text", "quote", "caption"):
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        txt = strip_md(txt)
        if kind == "bullet":
            run(p, "●  " + txt, size)
        elif kind == "sub":
            run(p, "      ○  " + txt, size - 1)
        elif kind == "quote":
            p.space_after = Pt(0)
            run(p, "      " + txt, size - 2, QUOTE, font=MONO)
        elif kind == "caption":
            p.alignment = PP_ALIGN.CENTER
            run(p, txt, size - 2, GREY)
        else:
            run(p, txt, size)
    return tb


def add_table(slide, rows, top):
    cells = [[c.strip() for c in r.strip("|").split("|")] for r in rows if not re.match(r"^\|[-\s|]+\|$", r)]
    nrows, ncols = len(cells), len(cells[0])
    row_h = Inches(0.42)
    shp = slide.shapes.add_table(nrows, ncols, BODY_BOX[0], top, BODY_BOX[2], row_h * nrows)
    tbl = shp.table
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = BODY_BOX[2] - Inches(2.2)
    for i, row in enumerate(cells):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            run(cell.text_frame.paragraphs[0], strip_md(val), 13, BLACK if i == 0 else BODY, bold=(i == 0))
    return row_h * nrows


def add_figure(slide, png, top, max_h):
    pic = slide.shapes.add_picture(str(png), 0, top)
    max_w = W - 2 * BODY_BOX[0]
    ratio = min(max_w / pic.width, max_h / pic.height)
    pic.width, pic.height = int(pic.width * ratio), int(pic.height * ratio)
    pic.left = int((W - pic.width) / 2)
    pic.top = top
    return pic.height


# ---------- build ----------
for s in slides:
    slide = prs.slides.add_slide(BLANK)
    t = s["title"]
    items = s["items"]
    if s["n"] in (1, 2):
        cover(slide, t, next((v for k, v in items if k == "text"), ""))
    elif t.startswith("[divider]"):
        divider(slide, t.replace("[divider]", "").strip())
    else:
        title(slide, re.sub(r"\s*\(figure\)$", "", t))
        fig = next((v for k, v in items if k == "figure"), None)
        tbl_rows = [v for k, v in items if k == "table"]
        x, y, w, h = BODY_BOX
        if fig:
            fh = add_figure(slide, FIGDIR / fig.replace(".svg", ".png"), y, Inches(3.0))
            add_body(slide, [(k, v) for k, v in items if k == "caption"], (x, y + fh + Inches(0.15), w, Inches(0.9)), 16)
        elif tbl_rows:
            th = add_table(slide, tbl_rows, y)
            rest = [(k, v) for k, v in items if k != "table"]
            add_body(slide, rest, (x, y + th + Inches(0.25), w, h - th - Inches(0.25)), 15)
        else:
            add_body(slide, items, BODY_BOX, body_size(items))
    if s["notes"]:
        slide.notes_slide.notes_text_frame.text = "\n".join(s["notes"])

prs.save(OUT)
print("slides", len(prs.slides), "->", OUT)
